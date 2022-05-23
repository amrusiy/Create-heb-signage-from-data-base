from pptx import Presentation
import pandas as pd
import os
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
import time
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import qrcode
import pyautogui as pg
import py_win_keyboard_layout
import win32gui,win32con
from barcode import Code128 #create barcode
from barcode.writer import ImageWriter
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE


data = pd.read_excel (r'F:\שילוטים\שילוטים חנות גדול\פה מכניסים הצעת מחיר.xlsx') #place "r" before the path string to address special character, such as '\'. Don't forget to put the file name at the end of the path + '.xlsx'

code_list = data["קוד פריט"].tolist()
name_of_product_list = data["תיאור פריט"].tolist()
link_list = data["קישור למפרט טכני נוסחא"].tolist()
test_link_list = data["כיתוב ליד QR נוסחא"].tolist()
price = data["מחיר ברוטו לאחר הנחה"].tolist()

num_of_row=len(code_list)
num_of_slide=int(num_of_row/4+1)
print("the number of  rows in the excel is: ",num_of_row)
print("The number of slide that we need is: ",num_of_slide)

##create new presantation
prs = Presentation()
layout = prs.slide_layouts[6]
## add a slide with the above layout
slide = prs.slides
slide = slide.add_slide(layout)
##change the prs size to 9:16
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

def search_and_replace(search_str, repl_str, input,pos):
    """"search and replace text in PowerPoint while preserving formatting"""
    #Useful Links ;)
    #https://stackoverflow.com/questions/37924808/python-pptx-power-point-find-and-replace-text-ctrl-h
    #https://stackoverflow.com/questions/45247042/how-to-keep-original-text-formatting-of-text-with-python-powerpoint
    prs = Presentation(input)
    count =0
    for slide in prs.slides:
        if count <= pos:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if(shape.text.find(search_str)) != -1:
                        text_frame = shape.text_frame
                        cur_text = text_frame.paragraphs[0].runs[0].text
                        new_text = cur_text.replace(str(search_str), str(repl_str))
                        text_frame.paragraphs[0].runs[0].text = new_text
        count=+count
    prs.save(name + '.pptx')
def fix_the_text(name,pos):
    search_and_replace('[','^',name+'.pptx',pos)
    search_and_replace(']','@',name+'.pptx',pos)
    search_and_replace('^',']',name+'.pptx',pos)
    search_and_replace('@','[',name+'.pptx',pos)
    search_and_replace('"',' "',name+'.pptx',pos)
    search_and_replace('(','^',name+'.pptx',pos)
    search_and_replace(')','@',name+'.pptx',pos)
    search_and_replace('^',')',name+'.pptx',pos)
    search_and_replace('@','(',name+'.pptx',pos)
    search_and_replace('\'[','[\'',name+'.pptx',pos)
    search_and_replace('+','+',name+'.pptx',pos)
    search_and_replace('יחידות','יח\'',name+'.pptx',pos)
    search_and_replace('[ ','[',name+'.pptx',pos)
    search_and_replace(' ]',']',name+'.pptx',pos)
    search_and_replace(']','] ',name+'.pptx',pos)
    search_and_replace('  ',' ',name+'.pptx',pos)
    search_and_replace('   ',' ',name+'.pptx',pos)
    search_and_replace('    ',' ',name+'.pptx',pos)
    search_and_replace('\'\'','',name+'.pptx',pos)
 #   letter = ('ט','ס','ד','S','T','ש','ק','מ','י','כ','ח')
  #  number = ('1','2','3','4','5','6','7','8','9','0')
  #  for let in letter:
  #      for num in number:
  #          no_sapce = num+let
  #          space = num + " " + let
  #          search_and_replace(no_sapce,space, name + '.pptx',pos)
def str_code(code):
    code=str(code)
    #code=code.replace(".0","")
    code=code.replace("b","")
    code=code.replace("'","")
    return code
def barcode(number = '1011673'):
    print(number)
    my_code = Code128(number,writer=ImageWriter())     # Now, let's create an object of EAN13 class and pass the number
    my_code.save("new_code")    # Our barcode is ready. Let's save it.
def create_qrcode(data):
    feature = qrcode.QRCode(version=1, box_size=30, border=0)
    feature.add_data(data)
    feature.make(fit=True)
    img_qr = qrcode.make = feature.make_image(fill_color="black", back_color="white")
    img_qr.save('qr_img.png')
    ##define counter that count the numver of sign that ready
##create user interface
def create_pattern_prs(x,product_code,product_name,qr_link,qr_text,product_price):
    ##create RECTANGLE pattern
    left = top = width = height = Inches(0.1)
    left=Inches(x+0.1)
    height = Inches(5.508110236)
    width = Inches(3.890866142)
    shape1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

    ##change the backround to white

    fill = shape1.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    ##change the line color to black
    fill_line = shape1.line
    fill_line.width = Inches(0.05)
    fill_line.color.rgb = RGBColor(0, 0, 0)
    #create text box
    left = Inches(0.13+x)
    if len(product_name)>60:
        top=Inches(0.2)
    elif len(product_name)>20 & len(product_name)<60:
        top = Inches(0.3)
    width = Inches(3.800866142)
    height = Inches(1.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    my_text = product_name
    run.text = my_text
    tf.word_wrap = True
    txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    txBox.text_frame.paragraphs[0].auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    font = run.font
    if len(run.text)<=25:
        font.size = Pt(38)
    elif len(run.text)<=50:
        font.size = Pt(32)
    elif  len(run.text)<=80 and len(run.text)>=50:
        font.size = Pt(26)
    else:
        font.size = Pt(22)
    ##add picture Frame for stickers
    img_frame_for_stickers='Frame for stickers.png'
    left = Inches(0.33+x)
    top = Inches(2.05)
    pic = slide.shapes.add_picture(img_frame_for_stickers, left, top)

    #add barcode image
    barcode(str(str_code(product_code)))
    barcode_img='new_code.png'
    left = Inches(0.8 + x)
    top = Inches(2.1)
    width = Inches(2.4)
    height = Inches(1.15)
    pic_barcode = slide.shapes.add_picture(barcode_img, left, top,width,height)

    #add white text to covert product code image
    left = Inches(x + 1.5)
    top = Inches(3)
    width = Inches(1.5)
    height = Inches(0.3)
    shape1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape1.shadow.inherit = False
    ##change the backround to white
    fill = shape1.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    line = shape1.line
    line.color.rgb = RGBColor(255, 255, 255)
    #add the product code
    left = Inches(x + 1.23)
    top = Inches(2.87)
    width = Inches(1.5)
    height = Inches(0.3)
    txBox_code = slide.shapes.add_textbox(left, top, width, height)
    txBox_code.line.fill.background()
    tf = txBox_code.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    my_text = product_code
    run.text = my_text
    tf.word_wrap = True
    txBox_code.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    font = run.font
    font.size = Pt(24)


    ##add text box for price
    left = Inches(x+1.6)
    top = Inches(3.25)
    width = Inches(1.3)
    height = Inches(0.2)
    txBox_qr = slide.shapes.add_textbox(left, top, width, height)
    tf_qr = txBox_qr.text_frame
    p_qr = tf_qr.paragraphs[0]
    run_qr = p_qr.add_run()
    run_qr.text = product_price
    tf_qr.word_wrap = True
    txBox_qr.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    font_qr = run_qr.font
    font_qr.size = Pt(18)

    ##add text box for "include taxs" near the price
    left = Inches(x+0.5)
    top = Inches(3.25)
    width = Inches(1.2)
    height = Inches(0.2)
    txBox_qr = slide.shapes.add_textbox(left, top, width, height)
    tf_qr = txBox_qr.text_frame
    p_qr = tf_qr.paragraphs[0]
    run_qr = p_qr.add_run()
    run_qr.text = "כולל מע\"מ"
    tf_qr.word_wrap = True
    txBox_qr.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    font_qr = run_qr.font
    font_qr.size = Pt(16)


    ##add text box for QR
    left = Inches(2.1+x)
    top=Inches(3.9)
    width = Inches(1.700866142)
    height = Inches(1)
    txBox_qr = slide.shapes.add_textbox(left, top, width, height)
    tf_qr = txBox_qr.text_frame
    p_qr = tf_qr.paragraphs[0]
    run_qr = p_qr.add_run()
    run_qr.text = qr_text
    tf_qr.word_wrap = True
    txBox_qr.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    font_qr = run_qr.font
    font_qr.size = Pt(22)
    #create function of QR image

    create_qrcode(qr_link)
    ##add QR image
    img_frame_for_stickers='qr_img.png'
    left = Inches(0.8+x)
    top = Inches(3.89)
    height = Inches(0.9)
    pic_frame_for_stickers = slide.shapes.add_picture(img_frame_for_stickers, left, top,height)
    ##add amrusi logo - image
    img_amrusi_logo ='amrusi logo.jpg'
    left = Inches(1.27+x)
    top = Inches(5.0)
    pic_img_amrusi_logo = slide.shapes.add_picture(img_amrusi_logo, left, top)

count = 0
pos = 0
for j in range(num_of_slide):
    x = 0
    for i in range(4):
        if count <= len(code_list)-1:
             ## information from excel
             product_code = code_list[count]
             if product_code > 100:
                 product_name = name_of_product_list[count]
                 product_code = str(product_code).replace(".0","")
                 product_name=str(product_name)
                 qr_link = link_list[count]
                 qr_text = test_link_list[count]
                 product_price = price[count]
                 product_price= str(round(product_price, 2))
                 product_price=str(product_price+"₪")
                 product_code = str(product_code).encode('utf-8', 'ignore')
                 print(product_price)
                 create_pattern_prs(x, product_code, product_name, qr_link, qr_text,product_price)
                 x = x + 3.870866142
                 count = count + 1
             else:
                 if pos == 0:
                    pos = j
                    break

    layout = prs.slide_layouts[6]
    slide = prs.slides
    slide = slide.add_slide(layout)
    print("slide number: ",j)


name = "שילוטים חנות גדול"
prs.save(name+'.pptx')
fix_the_text(name,pos)
py_win_keyboard_layout.change_foreground_window_keyboard_layout(0x04090409)
os.startfile(name+'.pptx')
time.sleep(3)
##open presantion on full size
hwnd = win32gui.GetForegroundWindow()
win32gui.ShowWindow(hwnd, win32con.SW_MAXIMIZE)
pg.click(500,800)

def change_text_right_to_eft(pos):
    for i in range(pos+1):
        pg.hotkey('ctrl','a')
        pg.hotkey('ctrl','a')
        pg.press('alt')
        py_win_keyboard_layout.change_foreground_window_keyboard_layout(0x040d)
        time.sleep(0.3)
        pg.press('h')
        time.sleep(0.1)
        pg.press('0')
        py_win_keyboard_layout.change_foreground_window_keyboard_layout(0x04090409)
        time.sleep(0.1)
        pg.press('esc')
        time.sleep(0.3)
        pg.press('down')
        time.sleep(0.4)
time.sleep(5)
change_text_right_to_eft(pos)

